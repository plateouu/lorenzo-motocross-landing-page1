import subprocess
import sys
import os
import platform

# --- PART 1: AUTO-INSTALLER ---
def install_library():
    print("Checking for required library (python-pptx)...")
    try:
        import pptx
    except ImportError:
        print("Library not found. Installing now...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
            print("Installation successful!")
        except Exception as e:
            print(f"Could not install automatically. Error: {e}")
            print("Please try running 'pip install python-pptx' in your terminal manually.")
            sys.exit(1)

# Run the installer before importing
install_library()

# --- PART 2: THE PRESENTATION GENERATOR ---
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_cited_imp():
    print("Generating presentation... Please wait.")
    prs = Presentation()

    def add_slide(title_text, bullets, search_query, citation_text):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        slide.shapes.title.text = title_text
        
        # Bullets
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True
        if bullets:
            tf.text = bullets[0]
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
        
        # Visual Box
        left = Inches(5.5); top = Inches(2.0); width = Inches(4.0); height = Inches(3.0)
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
        shape.text = f"STEP 1: Google Image Search for:\n'{search_query}'\nSTEP 2: Paste here."
        
        # Citation
        txBox = slide.shapes.add_textbox(left, top + height + Inches(0.1), width, Inches(0.5))
        p = txBox.text_frame.add_paragraph()
        p.text = f"Image Source: {citation_text}"
        p.font.size = Pt(10); p.font.italic = True; p.font.color.rgb = RGBColor(80, 80, 80)

    def add_section(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    # --- SLIDES ---
    add_section("The Creaking Axis: Digital Communities as Social Infrastructure", "AP Seminar IMP\n[Your Name]")
    
    add_slide("The Axis of Memory (Stimulus)", 
              ["Stimulus: Haruki Murakami, A Walk to Kobe (2002)", "Concept: Identity is tied to physical 'scenery'.", "Relevance: What happens when scenery is destroyed?"],
              "Haruki Murakami A Walk to Kobe book cover", "Murakami, 2002")
    
    add_slide("The Crisis of Connection", 
              ["U.S. Surgeon General (2023):", "- Lack of connection = 29% increased risk of premature death.", "- Equivalent to smoking 15 cigarettes/day."],
              "Surgeon General Advisory 2023 Social Connection Graph", "Office of the U.S. Surgeon General, 2023")
    
    add_section("Research Question", "To what extent could digital communities be considered meaningful social spaces for those who are isolated from physical spaces?")
    
    add_slide("Thesis & Economic Impact", 
              ["Thesis: Virtual communities are necessary 'Third Places'.", "Significance: Isolation costs Medicare $6.7 Billion/year (AARP)."],
              "AARP Public Policy Institute isolation cost chart", "AARP Public Policy Institute, 2017")
    
    add_slide("Why We Can't Just 'Go Outside'", 
              ["Danah Boyd (2014): Curfews restrict youth.", "Ben Goldfarb (2023): Sprawl creates 'Social Deserts'."],
              "Suburban sprawl aerial view cul de sac", "Goldfarb, 2023")
    
    add_slide("Digital Communities as 'Third Places'", 
              ["Steinkuehler (2006): MMOs are neutral grounds.", "Hampton (2011): Internet use = larger networks."],
              "World of Warcraft gameplay screenshot", "Steinkuehler & Williams, 2006")
    
    add_slide("The Tension: Surveillance Capitalism", 
              ["Zuboff (2019): Platforms are for profit, not connection.", "Turkle (2011): 'Illusion of Companionship'."],
              "The Age of Surveillance Capitalism book cover", "Zuboff, 2019")
    
    add_slide("Biological Limitations", 
              ["Cacioppo (2008): Text fails to down-regulate stress.", "The Trade-off: Imperfect connection > Toxic isolation."],
              "HPA axis stress response diagram", "Cacioppo, 2008 / Holt-Lunstad, 2023")
    
    add_slide("Solution: Digital Citizenship", 
              ["Mike Ribble (2015): Shift to 'Participation'.", "Action: Goal-based hybrid groups (Esports)."],
              "Mike Ribble 9 elements of digital citizenship infographic", "Ribble, 2015")
    
    add_slide("Limitations & Implications", 
              ["Implication: Legitimacy reduces stigma.", "Limitation: 'Dark Participation' (Bullying) requires facilitation."],
              "Cyberbullying vs Digital Mentorship illustration", "Quandt, 2015")
    
    add_slide("Conclusion", 
              ["Summary: Physical axis is broken; Digital is the lifeline.", "Call to Action: Fund 'Digital Placemaking'."],
              "Digital town square abstract concept art", "Concept: Goldfarb, 2023")

    # Save & Open
    script_dir = os.path.dirname(os.path.abspath(__file__))
    filename = 'IMP_Auto_Fixed.pptx'
    filepath = os.path.join(script_dir, filename)
    
    prs.save(filepath)
    print(f"\nSUCCESS! Created {filepath}")
    
    try:
        if platform.system() == 'Darwin': subprocess.call(('open', filepath))
        elif platform.system() == 'Windows': os.startfile(filepath)
        else: subprocess.call(('xdg-open', filepath))
    except Exception as e:
        print(f"Could not open file automatically: {e}")

if __name__ == "__main__":
    create_cited_imp()
